"""Document ingestion pipeline: extract text from PDF/DOCX/PPTX, chunk it."""

import os
import re
import logging
from dataclasses import dataclass, field
from typing import List, Optional

import pdfplumber

from rag.config import CHUNK_SIZE, CHUNK_OVERLAP, CHARS_PER_TOKEN, OCR_MIN_CHARS_PER_PAGE

logger = logging.getLogger(__name__)


# ---- Helpers for enhanced extraction ----

def _table_to_markdown(table: list) -> str:
    """Convert a 2D list (rows of cells) into a markdown table string."""
    if not table or not table[0]:
        return ""
    cleaned = []
    for row in table:
        cleaned.append([str(cell).strip() if cell else "" for cell in row])
    if not cleaned:
        return ""
    num_cols = len(cleaned[0])
    lines = []
    lines.append("| " + " | ".join(cleaned[0]) + " |")
    lines.append("| " + " | ".join(["---"] * num_cols) + " |")
    for row in cleaned[1:]:
        while len(row) < num_cols:
            row.append("")
        lines.append("| " + " | ".join(row[:num_cols]) + " |")
    return "\n".join(lines)


def _detect_section_header(line: str) -> bool:
    """Heuristic: detect if a line looks like a section heading."""
    stripped = line.strip()
    if not stripped or len(stripped) > 120:
        return False
    # Numbered section patterns: "1.", "1.1", "I.", "A.", "Capitulo 3"
    if re.match(r'^(?:\d+\.?\d*|[IVXLC]+\.|[A-Z]\.)\s+', stripped):
        return True
    # All-caps lines (3+ chars, no trailing period)
    if stripped.isupper() and len(stripped) >= 3 and not stripped.endswith('.'):
        return True
    # Title-case lines under 80 chars without trailing punctuation
    if stripped.istitle() and len(stripped) < 80 and stripped[-1] not in '.,:;!?':
        return True
    return False


def _split_into_segments(text: str) -> list:
    """Split text into segments: regular paragraphs and [TABLE]...[/TABLE] blocks."""
    segments = []
    parts = re.split(r'(\[TABLE\].*?\[/TABLE\])', text, flags=re.DOTALL)
    for part in parts:
        part = part.strip()
        if not part:
            continue
        if part.startswith("[TABLE]"):
            segments.append(part)
        else:
            paragraphs = re.split(r'\n\s*\n|\n', part)
            for p in paragraphs:
                p = p.strip()
                if p:
                    segments.append(p)
    return segments


def _prepend_section(section_header: str, chunk_text: str) -> str:
    """Prepend section context to a chunk if available."""
    if section_header and not chunk_text.startswith(section_header):
        return f"[Section: {section_header}]\n{chunk_text}"
    return chunk_text


@dataclass
class PageText:
    page_number: int
    text: str


@dataclass
class Chunk:
    doc_id: str
    filename: str
    page_number: int
    chunk_index: int
    text: str


def extract_text_pdf(filepath: str) -> List[PageText]:
    """Extract text from PDF using pdfplumber, preserving table structure."""
    pages = []
    try:
        with pdfplumber.open(filepath) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text() or ""

                # Extract tables and convert to markdown
                try:
                    tables = page.extract_tables() or []
                    if tables:
                        table_parts = []
                        for table in tables:
                            md = _table_to_markdown(table)
                            if md:
                                table_parts.append(md)
                        if table_parts:
                            tables_block = "\n\n[TABLE]\n" + "\n\n[TABLE]\n".join(table_parts) + "\n[/TABLE]\n"
                            text = text + "\n" + tables_block
                except Exception as te:
                    logger.debug(f"Table extraction failed on page {i+1}: {te}")

                if len(text.strip()) < OCR_MIN_CHARS_PER_PAGE:
                    ocr_text = _ocr_page(filepath, i)
                    if ocr_text:
                        text = ocr_text
                pages.append(PageText(page_number=i + 1, text=text))
    except Exception as e:
        logger.error(f"Error extracting PDF {filepath}: {e}")
    return pages


def _ocr_page(filepath: str, page_index: int) -> Optional[str]:
    """OCR a single page of a PDF. Returns None if OCR deps not available."""
    try:
        from pdf2image import convert_from_path
        import pytesseract

        images = convert_from_path(filepath, first_page=page_index + 1, last_page=page_index + 1)
        if images:
            text = pytesseract.image_to_string(images[0], lang="spa+eng")
            return text.strip()
    except ImportError:
        logger.debug("OCR dependencies (pytesseract/pdf2image) not available, skipping OCR")
    except Exception as e:
        logger.debug(f"OCR failed for page {page_index}: {e}")
    return None


def extract_text_docx(filepath: str) -> List[PageText]:
    """Extract text from DOCX, preserving tables and heading structure."""
    try:
        from docx import Document
        doc = Document(filepath)

        parts = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            # Detect headings via style name
            style_name = para.style.name if para.style else ""
            if "Heading" in style_name or "heading" in style_name:
                parts.append(f"## {text}")
            else:
                parts.append(text)

        # Extract tables as markdown
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(cells)
            md = _table_to_markdown(rows)
            if md:
                parts.append(f"\n[TABLE]\n{md}\n[/TABLE]\n")

        full_text = "\n".join(parts)
        blocks = _split_into_blocks(full_text, block_size=3000)
        return [PageText(page_number=i + 1, text=block) for i, block in enumerate(blocks)]
    except Exception as e:
        logger.error(f"Error extracting DOCX {filepath}: {e}")
        return []


def extract_text_pptx(filepath: str) -> List[PageText]:
    """Extract text from PPTX. Each slide becomes a page."""
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        pages = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)
            if texts:
                pages.append(PageText(page_number=i + 1, text="\n".join(texts)))
        return pages
    except Exception as e:
        logger.error(f"Error extracting PPTX {filepath}: {e}")
        return []


def extract_text(filepath: str) -> List[PageText]:
    """Extract text from a document based on file extension."""
    ext = os.path.splitext(filepath)[1].lower()
    if ext == ".pdf":
        return extract_text_pdf(filepath)
    elif ext == ".docx":
        return extract_text_docx(filepath)
    elif ext == ".pptx":
        return extract_text_pptx(filepath)
    else:
        logger.error(f"Unsupported file format: {ext}")
        return []


def chunk_text(pages: List[PageText], doc_id: str, filename: str) -> List[Chunk]:
    """Split page texts into overlapping chunks, preserving tables and section context."""
    chunks = []
    chunk_index = 0
    max_chars = CHUNK_SIZE * CHARS_PER_TOKEN  # ~2000 chars
    overlap_chars = CHUNK_OVERLAP * CHARS_PER_TOKEN  # ~200 chars

    for page in pages:
        text = page.text.strip()
        if not text:
            continue

        segments = _split_into_segments(text)
        current_section = ""
        current_chunk = ""

        for segment in segments:
            is_table = segment.startswith("[TABLE]")

            # Track section headers
            if not is_table and _detect_section_header(segment):
                current_section = segment.strip()

            # Would adding this segment exceed the limit?
            if len(current_chunk) + len(segment) + 1 > max_chars and current_chunk:
                # Flush current chunk
                chunk_final = _prepend_section(current_section, current_chunk.strip())
                chunks.append(Chunk(
                    doc_id=doc_id,
                    filename=filename,
                    page_number=page.page_number,
                    chunk_index=chunk_index,
                    text=chunk_final
                ))
                chunk_index += 1
                # Overlap — but not after tables (don't repeat table data)
                if overlap_chars > 0 and len(current_chunk) > overlap_chars and not is_table:
                    current_chunk = current_chunk[-overlap_chars:] + "\n" + segment
                else:
                    current_chunk = segment
            else:
                current_chunk = current_chunk + "\n" + segment if current_chunk else segment

            # If a table itself exceeds max_chars, emit it as its own chunk
            if is_table and len(current_chunk) > max_chars:
                chunk_final = _prepend_section(current_section, current_chunk.strip())
                chunks.append(Chunk(
                    doc_id=doc_id,
                    filename=filename,
                    page_number=page.page_number,
                    chunk_index=chunk_index,
                    text=chunk_final
                ))
                chunk_index += 1
                current_chunk = ""

        if current_chunk.strip():
            chunk_final = _prepend_section(current_section, current_chunk.strip())
            chunks.append(Chunk(
                doc_id=doc_id,
                filename=filename,
                page_number=page.page_number,
                chunk_index=chunk_index,
                text=chunk_final
            ))
            chunk_index += 1

    return chunks


def process_document(filepath: str, doc_id: str, filename: Optional[str] = None) -> List[Chunk]:
    """Full pipeline: extract text from document and chunk it."""
    if filename is None:
        filename = os.path.basename(filepath)

    logger.info(f"Processing document: {filename}")
    pages = extract_text(filepath)
    logger.info(f"Extracted {len(pages)} pages from {filename}")

    chunks = chunk_text(pages, doc_id, filename)
    logger.info(f"Created {len(chunks)} chunks from {filename}")

    return chunks


def get_page_count(filepath: str) -> int:
    """Get the number of pages in a document."""
    ext = os.path.splitext(filepath)[1].lower()
    try:
        if ext == ".pdf":
            with pdfplumber.open(filepath) as pdf:
                return len(pdf.pages)
        elif ext == ".docx":
            from docx import Document
            doc = Document(filepath)
            # Rough estimate: 1 page per 3000 chars
            text_len = sum(len(p.text) for p in doc.paragraphs)
            return max(1, text_len // 3000 + 1)
        elif ext == ".pptx":
            from pptx import Presentation
            return len(Presentation(filepath).slides)
    except Exception:
        pass
    return 0


def _split_into_blocks(text: str, block_size: int = 3000) -> List[str]:
    """Split text into blocks of approximately block_size characters."""
    if len(text) <= block_size:
        return [text]

    blocks = []
    start = 0
    while start < len(text):
        end = start + block_size
        if end < len(text):
            # Try to break at a paragraph or sentence boundary
            for sep in ["\n\n", "\n", ". ", " "]:
                idx = text.rfind(sep, start, end)
                if idx > start:
                    end = idx + len(sep)
                    break
        blocks.append(text[start:end])
        start = end
    return blocks
